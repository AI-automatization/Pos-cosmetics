"""
RAOS Admin Panel — Professional User Guide Presentation
Design: Classic Blue + Cyan accent, with real screenshots from raos.uz
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from PIL import Image
import os

# ═══════════════ DESIGN SYSTEM ═══════════════
NAVY = RGBColor(0x0F, 0x17, 0x2A)
PRIMARY = RGBColor(0x1E, 0x40, 0xAF)
BLUE = RGBColor(0x3B, 0x82, 0xF6)
CYAN = RGBColor(0x06, 0xB6, 0xD4)
GREEN = RGBColor(0x10, 0xB9, 0x81)
AMBER = RGBColor(0xF5, 0x9E, 0x0B)
RED = RGBColor(0xEF, 0x44, 0x44)
DARK = RGBColor(0x1E, 0x29, 0x3B)
GRAY = RGBColor(0x64, 0x74, 0x8B)
LIGHT_GRAY = RGBColor(0x94, 0xA3, 0xB8)
SLATE = RGBColor(0xCB, 0xD5, 0xE1)
LIGHT = RGBColor(0xF1, 0xF5, 0xF9)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BORDER = RGBColor(0xE2, 0xE8, 0xF0)

SCREENSHOTS = os.path.join(os.path.expanduser("~"), "Desktop", "RAOS_Screenshots")

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW = prs.slide_width
SH = prs.slide_height

def rect(s, l, t, w, h, fill=None, border=None, bw=Pt(0)):
    sh = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    sh.line.fill.background()
    if fill: sh.fill.solid(); sh.fill.fore_color.rgb = fill
    else: sh.fill.background()
    if border: sh.line.color.rgb = border; sh.line.width = bw
    return sh

def rrect(s, l, t, w, h, fill=None, border=None):
    sh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    sh.line.fill.background()
    if fill: sh.fill.solid(); sh.fill.fore_color.rgb = fill
    if border: sh.line.color.rgb = border; sh.line.width = Pt(1.5)
    return sh

def txt(s, l, t, w, h, text, sz=14, c=DARK, b=False, al=PP_ALIGN.LEFT, fn="Arial"):
    tb = s.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = text
    p.font.size = Pt(sz); p.font.color.rgb = c; p.font.bold = b
    p.font.name = fn; p.alignment = al
    return tb

def add_img(slide, img_file, left, top, max_w, max_h):
    """Add image maintaining aspect ratio within bounds."""
    path = os.path.join(SCREENSHOTS, img_file)
    if not os.path.exists(path): return
    with Image.open(path) as im:
        iw, ih = im.size
    ar = iw / ih
    if max_w / ar <= max_h:
        fw, fh = max_w, max_w / ar
    else:
        fh, fw = max_h, max_h * ar
    cx = left + (max_w - fw) // 2
    pic = slide.shapes.add_picture(path, int(cx), int(top), int(fw), int(fh))
    pic.line.color.rgb = BORDER; pic.line.width = Pt(1)

def header(s, title, subtitle=None):
    """Standard slide header bar."""
    rect(s, 0, 0, SW, Inches(1.0), fill=PRIMARY)
    rect(s, 0, 0, Inches(0.08), Inches(1.0), fill=CYAN)
    txt(s, Inches(0.5), Inches(0.18), Inches(10), Inches(0.5), title, sz=26, c=WHITE, b=True)
    if subtitle:
        txt(s, Inches(0.5), Inches(0.58), Inches(10), Inches(0.3), subtitle, sz=11, c=RGBColor(0x93, 0xC5, 0xFD))
    # Page number area
    rect(s, SW - Inches(0.8), Inches(0.3), Inches(0.5), Inches(0.4), fill=CYAN)

def screenshot_slide(title, subtitle, img_file):
    """Create a slide with header + full screenshot."""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    header(s, title, subtitle)
    rect(s, Inches(0.3), Inches(1.1), SW - Inches(0.6), SH - Inches(1.3), fill=LIGHT)
    add_img(s, img_file, Inches(0.4), Inches(1.15), SW - Inches(0.8), SH - Inches(1.4))
    return s

# ═══════════════════════════════════════════════════
# SLIDE 1: COVER
# ═══════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
rect(s, 0, 0, SW, SH, fill=NAVY)
# Accent lines
rect(s, 0, 0, Inches(0.08), SH, fill=CYAN)
rect(s, 0, Inches(6.5), SW, Inches(1.0), fill=RGBColor(0x0A, 0x0F, 0x1E))
# Decorative shapes
rect(s, Inches(10), Inches(0.5), Inches(3), Inches(0.04), fill=CYAN)
rect(s, Inches(0.5), Inches(6.3), Inches(5), Inches(0.04), fill=BLUE)

# Logo
logo = rrect(s, Inches(5.5), Inches(0.8), Inches(2.3), Inches(2.3), fill=WHITE)
txt(s, Inches(5.5), Inches(1.1), Inches(2.3), Inches(1.3), "RAOS", sz=52, c=PRIMARY, b=True, al=PP_ALIGN.CENTER)
txt(s, Inches(5.5), Inches(2.5), Inches(2.3), Inches(0.4), "v3.0", sz=16, c=GRAY, al=PP_ALIGN.CENTER)

# Title
txt(s, Inches(0.5), Inches(3.5), Inches(12.3), Inches(0.8), "RAOS — Admin Panel", sz=44, c=WHITE, b=True, al=PP_ALIGN.CENTER)
txt(s, Inches(0.5), Inches(4.3), Inches(12.3), Inches(0.5), "Foydalanish bo'yicha qo'llanma  |  Руководство пользователя", sz=20, c=RGBColor(0x93, 0xC5, 0xFD), al=PP_ALIGN.CENTER)
txt(s, Inches(0.5), Inches(5.2), Inches(12.3), Inches(0.5), "Retail & Asset Operating System  |  Multi-tenant  |  Offline-first  |  Kosmetika Savdosi", sz=13, c=RGBColor(0x64, 0x74, 0x8B), al=PP_ALIGN.CENTER)
txt(s, Inches(0.5), Inches(6.7), Inches(12.3), Inches(0.4), "raos.uz  |  2025-2026  |  Kosmetika Savdosi", sz=11, c=RGBColor(0x64, 0x74, 0x8B), al=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════
# SLIDE 2: TABLE OF CONTENTS
# ═══════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
header(s, "Mundarija  |  Содержание")

sections = [
    ("01", "Tizimga kirish", "Login sahifasi"),
    ("02", "Bosh sahifa", "Dashboard — KPI, grafik, top mahsulotlar"),
    ("03", "Katalog", "Mahsulotlar, kategoriyalar, yetkazuvchilar"),
    ("04", "Ombor", "Qoldiqlar, kirim, chiqim, harakatlar"),
    ("05", "Mijozlar", "Mijozlar bazasi va xarid tarixi"),
    ("06", "Xodimlar", "Rollar va vakolatlar boshqaruvi"),
    ("07", "Sotuvlar", "Buyurtmalar, qaytarishlar, smenalar"),
    ("08", "Hisobotlar", "Kunlik tushum, top mahsulotlar, filiallar"),
    ("09", "Analitika", "Biznes ko'rsatkichlari va trendlar"),
    ("10", "Moliya", "Nasiya, chegirmalar, xarajatlar"),
    ("11", "Ko'chmas mulk", "Ijaralar, shartnomalar, ROI"),
    ("12", "Sozlamalar", "Foydalanuvchilar, filiallar, printer, audit"),
]

for i, (num, title, desc) in enumerate(sections):
    col = 0 if i < 6 else 1
    row = i if i < 6 else i - 6
    x = Inches(0.8) + Inches(6.2) * col
    y = Inches(1.4) + Inches(0.95) * row

    # Number badge
    badge = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(0.55), Inches(0.55))
    badge.fill.solid(); badge.fill.fore_color.rgb = PRIMARY; badge.line.fill.background()
    tf = badge.text_frame; tf.paragraphs[0].text = num
    tf.paragraphs[0].font.size = Pt(14); tf.paragraphs[0].font.color.rgb = WHITE
    tf.paragraphs[0].font.bold = True; tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    txt(s, x + Inches(0.7), y - Inches(0.02), Inches(5), Inches(0.3), title, sz=15, c=DARK, b=True)
    txt(s, x + Inches(0.7), y + Inches(0.28), Inches(5), Inches(0.25), desc, sz=10, c=GRAY)

# ═══════════════════════════════════════════════════
# SCREENSHOT SLIDES — Login
# ═══════════════════════════════════════════════════
screenshot_slide("01  Tizimga kirish  |  Login", "Brauzerda raos.uz manzilini oching → Email va parolni kiriting → \"Kirish\" tugmasini bosing", "01_login.png")
screenshot_slide("01  Login — Maydonlar to'ldirilganda", "Email: owner@kosmetika.uz  |  Parol: ********  |  Kirish tugmasi", "02_login_filled.png")

# Dashboard
screenshot_slide("02  Bosh sahifa  |  Dashboard", "KPI kartochkalar, haftalik sotuv grafigi, filiallar daromadi, top mahsulotlar", "03_dashboard.png")

# Catalog
screenshot_slide("03  Katalog — Mahsulotlar", "26 ta mahsulot  |  SKU, narx, kategoriya, qoldiq  |  Qo'shish, tahrirlash, yorliq chop etish", "04_products.png")
screenshot_slide("03  Katalog — Kategoriyalar", "8 ta kategoriya  |  Yuz parvarishi, Soch, Parfyumeriya, Dekorativ, Tirnoq, Gigiena, Aksessuarlar", "05_categories.png")
screenshot_slide("03  Katalog — Yetkazib beruvchilar", "5 ta yetkazuvchi  |  L'Oréal, Nivea, Korean Beauty, P&G, Istanbul Cosmetics", "06_suppliers.png")

# Inventory
screenshot_slide("04  Ombor  |  Inventar", "Zaxira holati  |  Harakatlar tarixi  |  Nakladnoy  |  Kirim / Chiqim / Tester tugmalari", "07_inventory.png")

# Customers
screenshot_slide("05  Mijozlar  |  Customers", "12 ta mijoz  |  Ism, telefon, xarid tarixi  |  Filial bo'yicha filtr", "08_customers.png")

# Workers
screenshot_slide("06  Xodimlar  |  Workers", "8 ta xodim  |  Owner, Admin, Manager, Cashier, Warehouse  |  Filialga biriktirish", "09_workers.png")

# Reports
screenshot_slide("07  Hisobotlar  |  Reports", "Kunlik tushum  |  Top mahsulotlar  |  Filial solishtirish  |  Smena hisoboti", "10_reports.png")

# Analytics
screenshot_slide("08  Analitika  |  Analytics", "Umumiy daromad, buyurtmalar, o'rtacha chek  |  Sotuv trendi  |  ABC tahlil  |  Kassirlar", "11_analytics.png")

# Nasiya
screenshot_slide("09  Nasiya  |  Qarzlar", "Mijoz qarzlari  |  To'lov muddati  |  Qisman to'lov  |  Eslatma yuborish", "12_nasiya.png")

# Chegirma
screenshot_slide("09  Chegirmalar boshqaruvi", "Foizli va summaviy chegirmalar  |  Boshlanish/tugash sanasi  |  Faol/Kutilmoqda/Yakunlangan", "18_chegirma.png")

# Real Estate
screenshot_slide("10  Ko'chmas mulk  |  Real Estate", "Mulklar, ijara shartnomalar, to'lov jadvali, ROI hisoblash", "13_realestate.png")

# Promotions
screenshot_slide("09  Aksiyalar  |  Promotions", "Maxsus aksiya va promo-kodlar  |  Muddat belgilash  |  Statistika", "17_promotions.png")

# Exchange Rates
screenshot_slide("09  Valyuta kurslari", "USD, EUR, RUB kurslar  |  CBU dan avtomatik yangilanadi  |  Kunlik tarix", "19_exchange_rates.png")

# Tasks
screenshot_slide("12  Topshiriqlar  |  Tasks", "Jamoa vazifalari  |  Bajarilish holati  |  Prioritet  |  Mas'ul xodim", "16_tasks.png")

# Settings
screenshot_slide("12  Sozlamalar — Foydalanuvchilar", "Xodim qo'shish  |  Rol biriktirish (Owner/Admin/Manager/Cashier/Warehouse)  |  Parol tiklash", "14_settings_users.png")
screenshot_slide("12  Sozlamalar — Filiallar", "3 ta filial: Chilonzor, Mirzo Ulug'bek, Sergeli  |  Manzil  |  Xodimlar soni", "15_settings_branches.png")
screenshot_slide("12  Sozlamalar — Printer", "Chek printer sozlamalari  |  Printer turi  |  Test chek  |  Shablon tahrirlash", "23_printer.png")
screenshot_slide("12  Sozlamalar — Billing", "Tarif rejasi  |  To'lov holati  |  Obuna boshqaruvi", "20_billing.png")
screenshot_slide("12  Sozlamalar — Audit jurnal", "Barcha amallar logi  |  Kim, qachon, nima qildi  |  Sana bo'yicha filtr", "21_audit_log.png")

# Onboarding
screenshot_slide("Onboarding — Boshlang'ich sozlash", "Yangi foydalanuvchi uchun tizimni sozlash bosqichlari", "22_onboarding.png")

# ═══════════════════════════════════════════════════
# INFO SLIDE: ROLES & USERS
# ═══════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
header(s, "Foydalanuvchilar va rollar", "Tizimga kirish uchun login ma'lumotlari")

# Roles table
from pptx.util import Inches as I
roles_data = [
    ["Email", "Ism", "Rol", "Parol"],
    ["owner@kosmetika.uz", "Akbar Tursunov", "OWNER", "Demo2026!"],
    ["admin@kosmetika.uz", "Dilnoza Yusupova", "ADMIN", "Demo2026!"],
    ["manager.chilonzor@kosmetika.uz", "Jasur Toshmatov", "MANAGER", "Demo2026!"],
    ["manager.mirzo@kosmetika.uz", "Sardor Karimov", "MANAGER", "Demo2026!"],
    ["kassir.malika@kosmetika.uz", "Malika Rahimova", "CASHIER", "Demo2026!"],
    ["kassir.zulfiya@kosmetika.uz", "Zulfiya Nazarova", "CASHIER", "Demo2026!"],
    ["kassir.muhabbat@kosmetika.uz", "Muhabbat Aliyeva", "CASHIER", "Demo2026!"],
    ["ombor@kosmetika.uz", "Bobur Xolmatov", "WAREHOUSE", "Demo2026!"],
]

tbl = s.shapes.add_table(len(roles_data), 4, Inches(0.8), Inches(1.3), Inches(11.7), Inches(0.4) * len(roles_data))
table = tbl.table
for i, row in enumerate(roles_data):
    for j, val in enumerate(row):
        cell = table.cell(i, j); cell.text = val
        p = cell.text_frame.paragraphs[0]
        p.font.size = Pt(11); p.font.name = "Arial"
        if i == 0:
            cell.fill.solid(); cell.fill.fore_color.rgb = PRIMARY
            p.font.color.rgb = WHITE; p.font.bold = True
        else:
            cell.fill.solid()
            cell.fill.fore_color.rgb = WHITE if i % 2 == 1 else LIGHT
            p.font.color.rgb = DARK

# Role hierarchy
txt(s, Inches(0.8), Inches(5.2), Inches(11), Inches(0.4), "Rol ierarxiyasi:", sz=14, c=DARK, b=True)
roles_h = ["OWNER (5)", "ADMIN (4)", "MANAGER (3)", "CASHIER (2)", "WAREHOUSE", "VIEWER (1)"]
colors_h = [PRIMARY, BLUE, GREEN, AMBER, RGBColor(0x8B, 0x5C, 0xF6), LIGHT_GRAY]
for i, (role, color) in enumerate(zip(roles_h, colors_h)):
    x = Inches(0.8) + Inches(2.0) * i
    badge = rrect(s, x, Inches(5.7), Inches(1.7), Inches(0.5), fill=color)
    txt(s, x, Inches(5.77), Inches(1.7), Inches(0.35), role, sz=11, c=WHITE, b=True, al=PP_ALIGN.CENTER)
    if i < 5:
        txt(s, x + Inches(1.75), Inches(5.77), Inches(0.25), Inches(0.35), ">", sz=14, c=GRAY, b=True, al=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════
# CLOSING SLIDE
# ═══════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
rect(s, 0, 0, SW, SH, fill=NAVY)
rect(s, 0, 0, Inches(0.08), SH, fill=CYAN)
rect(s, 0, Inches(6.5), SW, Inches(1.0), fill=RGBColor(0x0A, 0x0F, 0x1E))

logo = rrect(s, Inches(5.5), Inches(1.2), Inches(2.3), Inches(2.3), fill=WHITE)
txt(s, Inches(5.5), Inches(1.5), Inches(2.3), Inches(1.3), "RAOS", sz=48, c=PRIMARY, b=True, al=PP_ALIGN.CENTER)
txt(s, Inches(5.5), Inches(2.7), Inches(2.3), Inches(0.4), "v3.0", sz=16, c=GRAY, al=PP_ALIGN.CENTER)

txt(s, Inches(0.5), Inches(3.9), Inches(12.3), Inches(0.6), "Savollar bo'lsa — administratorga murojaat qiling", sz=24, c=WHITE, al=PP_ALIGN.CENTER)
txt(s, Inches(0.5), Inches(4.6), Inches(12.3), Inches(0.4), "Telegram: @raos_support  |  Email: support@raos.uz  |  Web: raos.uz", sz=15, c=RGBColor(0x93, 0xC5, 0xFD), al=PP_ALIGN.CENTER)
txt(s, Inches(0.5), Inches(5.3), Inches(12.3), Inches(0.4), "Kosmetika Savdosi  |  Chilonzor  |  Mirzo Ulug'bek  |  Sergeli", sz=12, c=GRAY, al=PP_ALIGN.CENTER)
txt(s, Inches(0.5), Inches(6.7), Inches(12.3), Inches(0.4), "RAOS  2025-2026. Barcha huquqlar himoyalangan.", sz=11, c=RGBColor(0x64, 0x74, 0x8B), al=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════
# SAVE
# ═══════════════════════════════════════════════════
out = os.path.join(os.path.dirname(__file__), "RAOS_Admin_Panel_Guide.pptx")
prs.save(out)
print(f"Saved: {out}")
print(f"Slides: {len(prs.slides)}")

# Copy to Desktop
import shutil
desktop = os.path.join(os.path.expanduser("~"), "Desktop", "RAOS_Admin_Panel_Guide.pptx")
shutil.copy2(out, desktop)
print(f"Copied: {desktop}")
