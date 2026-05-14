"""
Insert real screenshots into RAOS Admin Panel PPTX presentation.
Each page gets its own slide with a header + screenshot image.
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import os

SCREENSHOTS_DIR = os.path.join(os.path.expanduser("~"), "Desktop", "RAOS_Screenshots")
OUTPUT = os.path.join(os.path.dirname(__file__), "RAOS_Admin_Panel_Guide.pptx")

# Load existing presentation
prs = Presentation(OUTPUT)

SLIDE_WIDTH = prs.slide_width
SLIDE_HEIGHT = prs.slide_height

# Design constants
PRIMARY = RGBColor(0x1E, 0x40, 0xAF)
ACCENT = RGBColor(0x06, 0xB6, 0xD4)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK = RGBColor(0x1F, 0x29, 0x37)
MEDIUM = RGBColor(0x64, 0x74, 0x8B)
BORDER = RGBColor(0xE2, 0xE8, 0xF0)
LIGHT_BG = RGBColor(0xF1, 0xF5, 0xF9)


def add_shape(slide, left, top, width, height, fill_color):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape


def add_text(slide, left, top, width, height, text, font_size=14, color=DARK, bold=False, alignment=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = "Arial"
    p.alignment = alignment
    return txBox


def add_screenshot_slide(prs, title, subtitle, screenshot_file, insert_before_idx=None):
    """Add a slide with header + screenshot image."""
    img_path = os.path.join(SCREENSHOTS_DIR, screenshot_file)
    if not os.path.exists(img_path):
        print(f"  SKIP: {screenshot_file} not found")
        return None

    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

    # Header bar
    add_shape(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.9), PRIMARY)
    add_shape(slide, Inches(0), Inches(0), Inches(0.1), Inches(0.9), ACCENT)

    # Title
    add_text(slide, Inches(0.5), Inches(0.15), Inches(10), Inches(0.5),
             title, font_size=22, color=WHITE, bold=True)

    # Subtitle
    if subtitle:
        add_text(slide, Inches(0.5), Inches(0.5), Inches(10), Inches(0.3),
                 subtitle, font_size=11, color=RGBColor(0x93, 0xC5, 0xFD))

    # Screenshot image — centered, with border effect
    # Image area: from y=1.0 to bottom with padding
    img_top = Inches(1.05)
    img_max_width = Inches(12.5)
    img_max_height = Inches(6.1)

    # Add light background behind image
    add_shape(slide, Inches(0.4), Inches(0.95), Inches(12.5), Inches(6.3), LIGHT_BG)

    # Add the screenshot
    from PIL import Image
    with Image.open(img_path) as im:
        img_w, img_h = im.size

    aspect = img_w / img_h
    # Fit within max bounds
    if img_max_width / aspect <= img_max_height:
        # Width is limiting
        final_w = img_max_width
        final_h = img_max_width / aspect
    else:
        # Height is limiting
        final_h = img_max_height
        final_w = img_max_height * aspect

    # Center horizontally
    img_left = (SLIDE_WIDTH - final_w) / 2

    pic = slide.shapes.add_picture(img_path, int(img_left), int(img_top), int(final_w), int(final_h))

    # Add thin border to image
    pic.line.color.rgb = BORDER
    pic.line.width = Pt(1)

    print(f"  Added: {title} ({screenshot_file})")
    return slide


# ═══════════════════════════════════════════════════
# Define screenshot slides to insert
# ═══════════════════════════════════════════════════
screenshot_slides = [
    # (title, subtitle, filename)
    # After slide 3 (Login) — insert login screenshots
    ("Login sahifasi — Ekran ko'rinishi", "Brauzerda ko'rinadigan tizimga kirish oynasi", "01_login.png"),
    ("Login — Maydonlar to'ldirilganda", "Email va parol kiritilgan holat", "02_login_filled.png"),

    # After Dashboard section
    ("Dashboard — Bosh sahifa", "Haftalik sotuv grafigi, top mahsulotlar, sidebar navigatsiya", "03_dashboard.png"),

    # Catalog
    ("Katalog — Mahsulotlar", "Tovarlar ro'yxati sahifasi", "04_products.png"),
    ("Katalog — Kategoriyalar", "Mahsulot kategoriyalari boshqaruvi", "05_categories.png"),
    ("Katalog — Yetkazib beruvchilar", "Yetkazib beruvchilar (postavshhiklar) ro'yxati", "06_suppliers.png"),

    # Inventory
    ("Ombor (Inventar)", "Zaxira holati, harakatlar tarixi, nakladnoy tablari", "07_inventory.png"),

    # Customers & Workers
    ("Mijozlar", "Mijozlar bazasi — ism, telefon, xarid tarixi", "08_customers.png"),
    ("Xodimlar", "Xodimlar ro'yxati va rol boshqaruvi", "09_workers.png"),

    # Reports & Analytics
    ("Hisobotlar", "Kunlik tushum, top mahsulotlar, filial solishtirish", "10_reports.png"),
    ("Analitika", "Biznes ko'rsatkichlari: daromad, buyurtmalar, o'rtacha chek", "11_analytics.png"),

    # Finance
    ("Nasiya (Qarzlar)", "Mijozlar qarz holati va to'lov tarixi", "12_nasiya.png"),
    ("Chegirmalar boshqaruvi", "Mahsulotlarga chegirma berish — foizli va summaviy", "18_chegirma.png"),

    # Real Estate
    ("Ko'chmas mulk", "Ijaralar, shartnomalar, to'lov jadvali", "13_realestate.png"),

    # Promotions & Exchange
    ("Aksiyalar", "Maxsus aksiya va promo-kodlar boshqaruvi", "17_promotions.png"),
    ("Valyuta kurslari", "USD, EUR, RUB avtomatik yangilanuvchi kurslar", "19_exchange_rates.png"),

    # Management
    ("Topshiriqlar (Tasks)", "Jamoa vazifalari va bajarilish holati", "16_tasks.png"),

    # Settings
    ("Sozlamalar — Foydalanuvchilar", "Foydalanuvchi qo'shish, rol biriktirish, parol tiklash", "14_settings_users.png"),
    ("Sozlamalar — Filiallar", "Filiallar boshqaruvi va manzillar", "15_settings_branches.png"),
    ("Sozlamalar — Printer", "Chek printer sozlamalari va test chek", "23_printer.png"),
    ("Sozlamalar — Billing", "Tarif rejasi va to'lov holati", "20_billing.png"),
    ("Sozlamalar — Audit jurnal", "Barcha amallar logi — kim, qachon, nima qildi", "21_audit_log.png"),

    # Onboarding
    ("Onboarding — Boshlang'ich sozlash", "Yangi foydalanuvchi uchun tizimni sozlash bosqichlari", "22_onboarding.png"),
]

print(f"Existing slides: {len(prs.slides)}")
print(f"Adding {len(screenshot_slides)} screenshot slides...\n")

for title, subtitle, filename in screenshot_slides:
    add_screenshot_slide(prs, title, subtitle, filename)

# Save
prs.save(OUTPUT)
print(f"\nPresentation updated: {OUTPUT}")
print(f"Total slides: {len(prs.slides)}")
